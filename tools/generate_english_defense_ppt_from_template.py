from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"d:\OneDrive - Queen Mary, University of London\Desktop\code\AndroidCode\Chinese_game")
TEMPLATE = ROOT / "doc" / "Copyright PDFs Ally.pptx"
OUTPUT = ROOT / "Android_Chinese_Words_Game_App_Defence_English_Template_Polished.pptx"

DB_DIAGRAM = ROOT / "sql" / "sqlite_database.png"
MASCOT = ROOT / "output" / "8a7c0663-f7b2-4937-9eb5-b4aeb9d07c89_no_bg.png"
COVER_IMAGE = ROOT / "output" / "b6db589c-7942-45dd-b976-168c70c0582a.jpg"

DARK = RGBColor(27, 74, 58)
GREEN = RGBColor(104, 148, 84)
GREEN_SOFT = RGBColor(232, 242, 228)
GREEN_PALE = RGBColor(244, 249, 241)
GOLD = RGBColor(185, 155, 92)
TEXT = RGBColor(40, 46, 43)
MUTED = RGBColor(96, 106, 100)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(201, 214, 194)


def get_placeholders(slide):
    placeholders = list(slide.placeholders)
    if len(placeholders) < 2:
        raise ValueError("Expected at least two placeholders on the slide.")
    return placeholders[0], placeholders[1]


def clear_frame(text_frame, margin=6):
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Pt(margin)
    text_frame.margin_right = Pt(margin)
    text_frame.margin_top = Pt(margin)
    text_frame.margin_bottom = Pt(margin)


def set_title(shape, text, size=24):
    text_frame = shape.text_frame
    clear_frame(text_frame, margin=0)
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = DARK


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    clear_frame(text_frame, margin=2)
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(shape, bullets, size=18, color=TEXT, spacing=7):
    text_frame = shape.text_frame
    clear_frame(text_frame, margin=4)
    for idx, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.space_after = Pt(spacing)
        paragraph.line_spacing = 1.12
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color


def add_card(slide, left, top, width, height, title, bullets=None, fill=WHITE,
             title_color=DARK, body_size=15):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE

    add_textbox(slide, left + Inches(0.12), top + Inches(0.08), width - Inches(0.24), Inches(0.35),
                title, size=17, bold=True, color=title_color)
    if bullets:
        body = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.48),
                                        width - Inches(0.24), height - Inches(0.56))
        add_bullets(body, bullets, size=body_size, spacing=5)
    return shape


def add_stat_card(slide, left, top, width, height, value, label):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = GREEN_PALE
    card.line.color.rgb = LINE
    add_textbox(slide, left, top + Inches(0.08), width, Inches(0.35), value, size=22, bold=True,
                color=GREEN, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.05), top + Inches(0.48), width - Inches(0.1), Inches(0.32),
                label, size=11, color=MUTED, align=PP_ALIGN.CENTER)


def add_section_label(slide, left, top, width, text):
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.28))
    pill.fill.solid()
    pill.fill.fore_color.rgb = GREEN_SOFT
    pill.line.fill.background()
    add_textbox(slide, left, top + Inches(0.01), width, Inches(0.2), text, size=10, bold=True,
                color=GREEN, align=PP_ALIGN.CENTER)


def add_process_step(slide, left, top, width, title, body):
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(1.05))
    outer.fill.solid()
    outer.fill.fore_color.rgb = WHITE
    outer.line.color.rgb = LINE
    add_textbox(slide, left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.25),
                title, size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.42), width - Inches(0.24), Inches(0.44),
                body, size=11, color=MUTED, align=PP_ALIGN.CENTER)


def delete_extra_slides(prs, keep_count):
    for idx in range(len(prs.slides) - 1, keep_count - 1, -1):
        slide_id = prs.slides._sldIdLst[idx]
        prs.part.drop_rel(slide_id.rId)
        del prs.slides._sldIdLst[idx]


def cover_slide(slide):
    title_shape, body_shape = get_placeholders(slide)
    set_title(title_shape, "Android Chinese Words Game App", size=25)
    body_shape.text = ""

    add_section_label(slide, Inches(0.48), Inches(1.6), Inches(1.55), "UNDERGRADUATE DEFENCE")
    add_textbox(
        slide,
        Inches(0.5),
        Inches(1.95),
        Inches(4.6),
        Inches(1.0),
        "Design and Implementation of a Gamified Native Android App for Chinese Learning",
        size=20,
        bold=True,
        color=TEXT,
    )
    add_textbox(
        slide,
        Inches(0.5),
        Inches(3.05),
        Inches(4.8),
        Inches(0.9),
        "A compact learning system that combines vocabulary recognition, pronunciation practice, and sentence reconstruction.",
        size=13,
        color=MUTED,
    )

    meta = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.1), Inches(4.1), Inches(0.85))
    meta.fill.solid()
    meta.fill.fore_color.rgb = GREEN_PALE
    meta.line.color.rgb = LINE
    add_textbox(slide, Inches(0.65), Inches(4.28), Inches(3.8), Inches(0.45),
                "Liu Weiliang | BUPT 2022213200 / QMUL 221166552 | 20 April 2026",
                size=11, color=DARK)

    slide.shapes.add_picture(str(COVER_IMAGE), Inches(5.55), Inches(1.15), width=Inches(3.65), height=Inches(3.65))
    add_textbox(slide, Inches(6.0), Inches(4.95), Inches(2.75), Inches(0.3),
                "Simple, visual, and learner-friendly", size=11, bold=True, color=GREEN,
                align=PP_ALIGN.CENTER)


def background_slide(slide):
    title_shape, body_shape = get_placeholders(slide)
    set_title(title_shape, "Project Background and Problem Statement")
    body_shape.text = ""

    add_card(
        slide, Inches(0.5), Inches(1.2), Inches(4.25), Inches(1.65),
        "Learning Context",
        [
            "Mobile learning supports short, repeatable practice outside the classroom.",
            "Chinese learners must coordinate characters, pinyin, tones, and word order at the same time.",
        ],
        fill=GREEN_PALE,
    )
    add_card(
        slide, Inches(0.5), Inches(3.0), Inches(4.25), Inches(1.65),
        "Gap in Existing Tools",
        [
            "Many apps emphasise memorisation or display rather than active interaction and feedback.",
            "Pronunciation and sentence-level practice are often weak or disconnected from vocabulary work.",
        ],
    )
    add_card(
        slide, Inches(5.0), Inches(1.55), Inches(4.2), Inches(2.35),
        "Research Question",
        [
            "How can a native Android app use gamified interaction to improve vocabulary recognition, pronunciation training, and sentence understanding for beginner and intermediate Chinese learners?",
        ],
        fill=WHITE,
        body_size=14,
    )
    slide.shapes.add_picture(str(MASCOT), Inches(7.55), Inches(4.0), width=Inches(1.25), height=Inches(1.25))


def aim_slide(slide):
    title_shape, body_shape = get_placeholders(slide)
    set_title(title_shape, "Aim, Objectives and Contributions")
    body_shape.text = ""

    add_card(
        slide, Inches(0.5), Inches(1.25), Inches(4.15), Inches(3.4),
        "Project Aim",
        [
            "Design and implement a native Android prototype for Chinese learning.",
            "Support beginner and intermediate learners in a low-friction mobile setting.",
            "Combine learning tasks with progress, scores, and motivation mechanics.",
        ],
        fill=GREEN_PALE,
    )
    add_card(
        slide, Inches(4.9), Inches(1.25), Inches(4.3), Inches(3.4),
        "Main Contributions",
        [
            "Three integrated mini-games in one coherent learning flow.",
            "A practical Java/XML + SQLite architecture for an educational app.",
            "JSON-based question bank expansion without changing the core gameplay logic.",
            "A prototype route that connects UI, persistence, NLP support, and speech evaluation.",
        ],
        fill=WHITE,
    )


def overview_slide(slide):
    title_shape, body_shape = get_placeholders(slide)
    set_title(title_shape, "System Design Overview")
    body_shape.text = ""

    add_textbox(slide, Inches(0.55), Inches(1.18), Inches(8.6), Inches(0.38),
                "The app turns a learning session into a simple closed loop from entry to feedback.", size=13, color=MUTED)

    steps = [
        ("1. Sign In", "User login, profile, and personal progress entry."),
        ("2. Select Mode", "Choose game type and difficulty level."),
        ("3. Complete Session", "Play a 10-question training round."),
        ("4. Get Feedback", "Receive score, accuracy, and result details."),
    ]
    x = Inches(0.55)
    for title, body in steps:
        add_process_step(slide, x, Inches(1.85), Inches(2.05), title, body)
        x += Inches(2.18)

    add_card(
        slide, Inches(0.55), Inches(3.25), Inches(2.7), Inches(1.45),
        "Target Users",
        [
            "Beginner and intermediate Chinese learners.",
            "Users who benefit from short, repeatable mobile practice.",
        ],
        body_size=12,
    )
    add_card(
        slide, Inches(3.45), Inches(3.25), Inches(2.7), Inches(1.45),
        "Functional Scope",
        [
            "Account management, game flow, local records, and achievements.",
            "Three game types with difficulty-based question selection.",
        ],
        body_size=12,
        fill=GREEN_PALE,
    )
    add_card(
        slide, Inches(6.35), Inches(3.25), Inches(2.7), Inches(1.45),
        "Non-Functional Goals",
        [
            "Maintainable structure, responsive UI, and practical offline persistence.",
            "A prototype suited to extension rather than full product launch.",
        ],
        body_size=12,
    )


def architecture_slide(slide):
    title_shape, body_shape = get_placeholders(slide)
    set_title(title_shape, "Architecture and Technology Stack")
    body_shape.text = ""

    add_card(
        slide, Inches(0.55), Inches(1.25), Inches(4.1), Inches(3.5),
        "Architecture Route",
        [
            "The project moved from an early Unity/C# idea to a native Android implementation.",
            "Activity-based UI handles navigation, gameplay flow, and user interaction.",
            "A shared GameActivity coordinates the three game modes and session state.",
            "DAO-based persistence reduces SQL coupling inside interface code.",
        ],
        fill=GREEN_PALE,
    )
    add_card(
        slide, Inches(4.95), Inches(1.25), Inches(4.15), Inches(3.5),
        "Technology Stack",
        [
            "Java 11 + XML for Android application development.",
            "SQLite with DaoFactory and DAO implementations for local storage.",
            "HanLP and pinyin4j for Chinese text processing.",
            "iFlyTek-related modules for pronunciation evaluation support.",
        ],
    )


def modules_slide(slide):
    title_shape, body_shape = get_placeholders(slide)
    set_title(title_shape, "Core Learning Modules")
    body_shape.text = ""

    add_card(
        slide, Inches(0.45), Inches(1.35), Inches(2.8), Inches(3.2),
        "Character Matching",
        [
            "Choose the correct Chinese word from pinyin-guided options.",
            "Designed to strengthen the link between sound and written form.",
        ],
        fill=GREEN_PALE,
        body_size=13,
    )
    add_card(
        slide, Inches(3.55), Inches(1.35), Inches(2.8), Inches(3.2),
        "Pronunciation Quiz",
        [
            "Speak the target word and receive score-oriented pronunciation feedback.",
            "Turns oral practice into a repeatable game mechanic.",
        ],
        fill=WHITE,
        body_size=13,
    )
    add_card(
        slide, Inches(6.65), Inches(1.35), Inches(2.45), Inches(3.2),
        "Word Puzzle",
        [
            "Reorder short sentences or fill blanks in longer ones.",
            "Uses segmented sentence data to support question generation.",
        ],
        fill=GREEN_PALE,
        body_size=13,
    )


def data_slide(slide):
    title_shape, body_shape = get_placeholders(slide)
    set_title(title_shape, "Data Model and Local Persistence")
    body_shape.text = ""

    slide.shapes.add_picture(str(DB_DIAGRAM), Inches(0.6), Inches(1.22), height=Inches(3.45))
    caption = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.45), Inches(1.55), Inches(2.55), Inches(2.65))
    caption.fill.solid()
    caption.fill.fore_color.rgb = GREEN_PALE
    caption.line.color.rgb = LINE
    add_textbox(slide, Inches(6.62), Inches(1.75), Inches(2.2), Inches(0.3),
                "Why this matters", size=16, bold=True, color=DARK)
    add_textbox(
        slide, Inches(6.62), Inches(2.15), Inches(2.15), Inches(1.65),
        "The SQLite schema stores users, achievements, game scores, detailed question logs, sentences, segmented tokens, and question banks. JSON assets can be imported into the local database to expand content cleanly.",
        size=12, color=TEXT,
    )


def results_slide(slide):
    title_shape, body_shape = get_placeholders(slide)
    set_title(title_shape, "Current Implementation Status")
    body_shape.text = ""

    add_stat_card(slide, Inches(0.65), Inches(1.3), Inches(1.95), Inches(0.95), "3", "game modes")
    add_stat_card(slide, Inches(2.95), Inches(1.3), Inches(1.95), Inches(0.95), "3", "difficulty levels")
    add_stat_card(slide, Inches(5.25), Inches(1.3), Inches(1.95), Inches(0.95), "10", "questions per session")

    add_card(
        slide, Inches(0.6), Inches(2.55), Inches(5.15), Inches(2.05),
        "What is already working",
        [
            "Login, profile, game selection, and session-level game flow are implemented.",
            "All three game branches write results into the local database.",
            "Records, achievements, and progress data are already integrated into the prototype.",
        ],
        fill=WHITE,
        body_size=14,
    )
    add_card(
        slide, Inches(5.95), Inches(2.55), Inches(3.05), Inches(2.05),
        "Engineering Strengths",
        [
            "Shared gameplay controller.",
            "Expandable question-bank pipeline.",
            "Clear local persistence structure.",
        ],
        fill=GREEN_PALE,
        body_size=14,
    )


def future_slide(slide):
    title_shape, body_shape = get_placeholders(slide)
    set_title(title_shape, "Limitations and Future Work")
    body_shape.text = ""

    add_card(
        slide, Inches(0.6), Inches(1.3), Inches(4.15), Inches(3.2),
        "Current Limitations",
        [
            "Speech evaluation remains sensitive to pronunciation variation, background noise, and device conditions.",
            "Formal usability testing and broader performance validation are still incomplete.",
            "The interface is usable but not yet fully polished for final delivery quality.",
        ],
        fill=WHITE,
    )
    add_card(
        slide, Inches(5.0), Inches(1.3), Inches(4.1), Inches(3.2),
        "Next Steps",
        [
            "Stabilise the pronunciation-evaluation workflow and improve error handling.",
            "Expand the datasets and strengthen quality control for question generation.",
            "Add structured user evaluation, compatibility checks, and more refined visual feedback.",
        ],
        fill=GREEN_PALE,
    )


def closing_slide(slide):
    title_shape, body_shape = get_placeholders(slide)
    set_title(title_shape, "Conclusion and Q&A")
    body_shape.text = ""

    slide.shapes.add_picture(str(MASCOT), Inches(0.9), Inches(1.55), width=Inches(2.2), height=Inches(2.2))
    add_textbox(slide, Inches(3.5), Inches(1.7), Inches(5.2), Inches(0.65),
                "A feasible Android prototype for gamified Chinese learning", size=20, bold=True, color=DARK)
    add_textbox(
        slide, Inches(3.5), Inches(2.45), Inches(5.0), Inches(1.1),
        "The project integrates vocabulary recognition, pronunciation practice, sentence reconstruction, and local progress tracking in one coherent mobile system.",
        size=14,
        color=TEXT,
    )
    qa = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(3.85), Inches(2.7), Inches(0.72))
    qa.fill.solid()
    qa.fill.fore_color.rgb = GREEN_PALE
    qa.line.color.rgb = LINE
    add_textbox(slide, Inches(3.5), Inches(4.02), Inches(2.7), Inches(0.3),
                "Thank you. Questions are welcome.", size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation(str(TEMPLATE))
    delete_extra_slides(prs, 10)

    builders = [
        cover_slide,
        background_slide,
        aim_slide,
        overview_slide,
        architecture_slide,
        modules_slide,
        data_slide,
        results_slide,
        future_slide,
        closing_slide,
    ]

    for slide, builder in zip(prs.slides, builders):
        builder(slide)

    prs.save(str(OUTPUT))
    print(f"Generated: {OUTPUT}")


if __name__ == "__main__":
    main()
