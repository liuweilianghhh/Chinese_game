from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt


ROOT = Path(r"d:\OneDrive - Queen Mary, University of London\Desktop\code\AndroidCode\Chinese_game")
OUTPUT = ROOT / "Android_Chinese_Words_Game_App_答辩PPT_自动生成.pptx"
MASCOT = ROOT / "output" / "e62462db-040a-45f0-ade7-7e52858c23e0.png"
ER_DIAGRAM = ROOT / "ER图.jpg"


GREEN_DARK = RGBColor(18, 76, 53)
GREEN = RGBColor(70, 130, 86)
GREEN_LIGHT = RGBColor(220, 242, 224)
MINT = RGBColor(239, 248, 242)
CREAM = RGBColor(249, 247, 239)
SAND = RGBColor(230, 223, 206)
TEXT = RGBColor(36, 49, 43)
TEXT_LIGHT = RGBColor(95, 112, 103)
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(216, 180, 88)
RED = RGBColor(176, 72, 72)


def set_east_asian_font(run, font_name="Microsoft YaHei"):
    run.font.name = font_name
    rpr = run._r.get_or_add_rPr()
    rpr.set(qn("a:ea"), font_name)
    rpr.set(qn("a:latin"), font_name)
    rpr.set(qn("a:cs"), font_name)


def add_textbox(slide, left, top, width, height, text="", font_size=20, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei",
                fill=None, line_color=None, radius_shape=None, margin=0.15):
    if radius_shape:
        shape = slide.shapes.add_shape(radius_shape, left, top, width, height)
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
        else:
            shape.fill.background()
        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background()
        text_frame = shape.text_frame
    else:
        shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = shape.text_frame

    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Cm(margin)
    text_frame.margin_right = Cm(margin)
    text_frame.margin_top = Cm(margin)
    text_frame.margin_bottom = Cm(margin)
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_east_asian_font(run, font_name)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = align
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_paragraphs(shape, items, font_size=20, color=TEXT, bullet=True, level=0, spacing=4):
    tf = shape.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = level
        p.bullet = bullet
        p.space_after = Pt(spacing)
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            set_east_asian_font(run)
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP


def add_background(slide, color=CREAM):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Cm(0.55)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = GREEN_DARK
    slide.shapes[-1].line.fill.background()
    circle = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, prs.slide_width - Cm(6.8), -Cm(1.8), Cm(7.5), Cm(7.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN_LIGHT
    circle.line.fill.background()
    circle.fill.transparency = 0.15
    circle2 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, -Cm(1.8), prs.slide_height - Cm(3.5), Cm(4.2), Cm(4.2)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = SAND
    circle2.line.fill.background()
    circle2.fill.transparency = 0.1


def add_slide_title(slide, title, subtitle=None):
    add_textbox(slide, Cm(0.9), Cm(0.8), Cm(18), Cm(1.3), title, 24, True, GREEN_DARK)
    if subtitle:
        add_textbox(slide, Cm(0.95), Cm(2.0), Cm(20), Cm(0.8), subtitle, 11, False, TEXT_LIGHT)


def add_footer(slide, number):
    add_textbox(
        slide, prs.slide_width - Cm(1.7), prs.slide_height - Cm(0.8), Cm(1.0), Cm(0.5),
        str(number), 10, False, TEXT_LIGHT, PP_ALIGN.RIGHT
    )


def add_card(slide, left, top, width, height, title, body_items,
             fill=WHITE, title_fill=GREEN_DARK, title_color=WHITE):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = SAND
    head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Cm(1.0))
    head.fill.solid()
    head.fill.fore_color.rgb = title_fill
    head.line.fill.background()
    add_textbox(slide, left + Cm(0.2), top + Cm(0.1), width - Cm(0.4), Cm(0.7), title, 15, True, title_color)
    body = add_textbox(slide, left + Cm(0.25), top + Cm(1.1), width - Cm(0.5), height - Cm(1.25))
    add_paragraphs(body, body_items, font_size=13, color=TEXT)
    return card


def fit_image(slide, image_path, left, top, width=None, height=None):
    if width is not None and height is not None:
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    elif width is not None:
        slide.shapes.add_picture(str(image_path), left, top, width=width)
    else:
        slide.shapes.add_picture(str(image_path), left, top, height=height)


def make_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = GREEN_DARK

    left_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Cm(18.5), prs.slide_height)
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = GREEN_DARK
    left_panel.line.fill.background()

    right_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(18.5), 0, prs.slide_width - Cm(18.5), prs.slide_height)
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = CREAM
    right_panel.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(20.4), Cm(0.9), Cm(8.2), Cm(8.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN_LIGHT
    accent.fill.transparency = 0.08
    accent.line.fill.background()

    add_textbox(slide, Cm(1.1), Cm(1.8), Cm(15.5), Cm(1.0), "本科毕业设计答辩", 16, False, RGBColor(216, 238, 224))
    add_textbox(slide, Cm(1.1), Cm(3.0), Cm(16.0), Cm(3.0), "Android Chinese Words\nGame App", 24, True, WHITE)
    add_textbox(slide, Cm(1.1), Cm(6.3), Cm(15.8), Cm(1.2), "基于原生 Android 的中文词汇游戏学习应用设计与实现", 16, False, RGBColor(225, 239, 229))
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.1), Cm(8.0), Cm(5.0), Cm(1.0))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    badge.line.fill.background()
    add_textbox(slide, Cm(1.35), Cm(8.16), Cm(4.5), Cm(0.6), "项目代号：CatLingo", 13, True, GREEN_DARK)

    add_textbox(slide, Cm(1.15), Cm(11.0), Cm(7.8), Cm(0.8), "学生：刘维良", 13, False, WHITE)
    add_textbox(slide, Cm(1.15), Cm(11.8), Cm(7.8), Cm(0.8), "学号：BUPT 2022213200 / QMUL 221166552", 12, False, RGBColor(220, 231, 225))
    add_textbox(slide, Cm(1.15), Cm(12.6), Cm(7.8), Cm(0.8), f"日期：{date(2026, 4, 20).strftime('%Y年%m月%d日')}", 12, False, RGBColor(220, 231, 225))

    fit_image(slide, MASCOT, Cm(20.7), Cm(2.0), width=Cm(6.6), height=Cm(6.6))
    add_textbox(slide, Cm(19.9), Cm(9.5), Cm(7.6), Cm(2.2), "面向中文初学者与中级学习者\n融合识字、发音与句子重组训练", 15, True, GREEN_DARK, PP_ALIGN.CENTER)
    return slide


def make_background_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_slide_title(slide, "1. 研究背景与问题定义", "为什么要做一个面向中文学习者的游戏化 Android 应用")

    add_card(slide, Cm(1.0), Cm(3.0), Cm(8.5), Cm(7.2), "研究背景", [
        "移动学习碎片化、即时反馈强，已成为语言学习的重要补充方式。",
        "现有中文学习工具多偏向字词展示或机械记忆，互动性与持续激励不足。",
        "中文学习同时涉及汉字识别、拼音匹配、语音声调和句序理解，练习门槛较高。",
    ], fill=WHITE)
    add_card(slide, Cm(10.0), Cm(3.0), Cm(8.5), Cm(7.2), "核心问题", [
        "如何在 Android 端把词汇识别、发音练习与句子理解整合到同一学习系统中。",
        "如何通过积分、成就、难度分级与进度记录提升重复练习意愿。",
        "如何在本地持久化、语音评测与题库管理之间维持可维护的技术架构。",
    ], fill=WHITE, title_fill=GREEN)
    add_textbox(slide, Cm(19.3), Cm(3.4), Cm(8.3), Cm(7.0),
                "论文中的中心研究问题：\n\n“如何设计并实现一个基于 Android 的中文学习应用，以游戏化交互提升初中级学习者的词汇识别、发音训练和句子理解能力，并保持技术上可维护、教育上有意义？”",
                15, False, TEXT, align=PP_ALIGN.LEFT, fill=GREEN_LIGHT,
                radius_shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, margin=0.35)
    add_footer(slide, 2)
    return slide


def make_objective_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_slide_title(slide, "2. 研究目标与论文贡献", "从项目目标、功能边界到论文的主要贡献点")

    add_card(slide, Cm(1.0), Cm(3.0), Cm(8.7), Cm(8.3), "项目目标", [
        "设计并实现一款面向中文初中级学习者的游戏化 Android 应用。",
        "围绕汉字识别、发音练习、词汇记忆和句序重组构建三类小游戏。",
        "支持账户管理、难度选择、成绩记录、成就展示与学习进度追踪。",
    ], fill=WHITE)

    add_card(slide, Cm(10.1), Cm(3.0), Cm(8.7), Cm(8.3), "范围界定", [
        "聚焦本地化学习原型，不覆盖完整课程体系、云端同步或对话式 AI 教学。",
        "主要针对基础词汇、拼音对应、发音反馈和简单句子重构任务。",
        "目标是验证一条可扩展的教育应用技术路线，而非商业化成品。",
    ], fill=WHITE, title_fill=GREEN)

    contrib = add_textbox(slide, Cm(19.2), Cm(3.0), Cm(8.5), Cm(8.3), fill=MINT,
                          radius_shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, text="")
    contrib.text_frame.clear()
    title_p = contrib.text_frame.paragraphs[0]
    r = title_p.add_run()
    r.text = "论文贡献"
    set_east_asian_font(r)
    r.font.bold = True
    r.font.size = Pt(16)
    r.font.color.rgb = GREEN_DARK
    title_p.space_after = Pt(8)
    for line in [
        "整合字词匹配、发音测验、句子谜题三类任务，形成统一学习入口。",
        "将游戏化机制与学习目标绑定，而不是只做表面奖励。",
        "展示了 Java + XML + SQLite + NLP + 语音评测的原生 Android 实现方案。",
        "提供 JSON 到 SQLite 的题库加载与可持续扩展的数据管线。",
    ]:
        p = contrib.text_frame.add_paragraph()
        p.text = line
        p.bullet = True
        p.space_after = Pt(4)
        for run in p.runs:
            set_east_asian_font(run)
            run.font.size = Pt(13)
            run.font.color.rgb = TEXT
    add_footer(slide, 3)
    return slide


def make_requirements_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_slide_title(slide, "3. 系统需求与总体方案", "功能需求、非功能需求与学习场景")

    add_card(slide, Cm(1.0), Cm(3.0), Cm(8.4), Cm(4.2), "目标用户", [
        "中文初学者与中级学习者。",
        "适合短时、重复、低门槛的移动端练习场景。",
        "希望通过可视化进度和即时反馈建立持续学习习惯。",
    ], fill=WHITE)
    add_card(slide, Cm(1.0), Cm(7.6), Cm(8.4), Cm(4.0), "功能需求", [
        "登录注册、个人资料、密码修改。",
        "游戏选择、难度分级、成绩记录、成就系统。",
        "三类核心小游戏与题库加载机制。",
    ], fill=WHITE, title_fill=GREEN)

    add_card(slide, Cm(10.0), Cm(3.0), Cm(8.6), Cm(4.2), "非功能需求", [
        "界面清晰、操作轻量，适配一般 Android 设备。",
        "响应及时、结构可维护、UI 与数据访问解耦。",
        "支持原型阶段的稳定运行与后续扩展。",
    ], fill=WHITE)
    add_card(slide, Cm(10.0), Cm(7.6), Cm(8.6), Cm(4.0), "伦理与风险", [
        "本地存储用户信息，需关注隐私保护。",
        "语音评测存在口音偏差，不应被呈现为绝对判断。",
        "应明确说明第三方语音服务的处理过程与边界。",
    ], fill=WHITE, title_fill=GREEN)

    process = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(19.1), Cm(3.2), Cm(8.3), Cm(8.4))
    process.fill.solid()
    process.fill.fore_color.rgb = WHITE
    process.line.color.rgb = SAND
    add_textbox(slide, Cm(19.45), Cm(3.45), Cm(7.6), Cm(0.8), "学习闭环设计", 16, True, GREEN_DARK)
    steps = ["用户登录", "选择游戏与难度", "完成 10 题练习", "即时反馈与得分", "记录进度并解锁成就"]
    y = 4.5
    for i, step in enumerate(steps, start=1):
        bubble = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(19.5), Cm(y), Cm(0.8), Cm(0.8))
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = GREEN
        bubble.line.fill.background()
        add_textbox(slide, Cm(19.5), Cm(y + 0.02), Cm(0.8), Cm(0.6), str(i), 12, True, WHITE, PP_ALIGN.CENTER)
        add_textbox(slide, Cm(20.5), Cm(y - 0.02), Cm(5.8), Cm(0.7), step, 14, False, TEXT)
        if i < len(steps):
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(19.9), Cm(y + 0.8), Cm(19.9), Cm(y + 1.1))
            line.line.color.rgb = GREEN
        y += 1.35
    add_footer(slide, 4)
    return slide


def make_architecture_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_slide_title(slide, "4. 系统架构与技术路线", "从 Unity 规划迁移到原生 Android 的实现路径")

    timeline = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.0), Cm(3.0), Cm(26.8), Cm(2.1))
    timeline.fill.solid()
    timeline.fill.fore_color.rgb = WHITE
    timeline.line.color.rgb = SAND
    add_textbox(slide, Cm(1.3), Cm(3.25), Cm(26.0), Cm(0.6),
                "技术演进：Unity + C# 初始方案  →  原生 Android（Java 11 + XML）  →  SQLite + HanLP + pinyin4j + 语音评测模块",
                16, True, GREEN_DARK)
    add_textbox(slide, Cm(1.3), Cm(4.05), Cm(25.8), Cm(0.5),
                "迁移原因：更适合表单、导航、数据库和轻量教育交互，维护成本更低，和项目实际需求更匹配。", 12, False, TEXT_LIGHT)

    layers = [
        ("表示层", "MainActivity、game_choice、GameActivity、个人中心与成就页面"),
        ("控制层", "负责用户交互、模式切换、题目推进与页面跳转"),
        ("数据层", "SQLite + DAO + DataManager，管理用户、成绩、题库和进度"),
        ("外部服务层", "HanLP 分词与讯飞相关发音评测模块"),
    ]
    y = 6.0
    colors = [GREEN_DARK, GREEN, RGBColor(95, 143, 103), RGBColor(150, 173, 120)]
    for idx, (title, body) in enumerate(layers):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(2.2), Cm(y), Cm(13.0), Cm(1.35))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[idx]
        box.line.fill.background()
        add_textbox(slide, Cm(2.6), Cm(y + 0.14), Cm(2.3), Cm(0.8), title, 16, True, WHITE)
        add_textbox(slide, Cm(5.0), Cm(y + 0.14), Cm(9.6), Cm(0.8), body, 12, False, WHITE)
        y += 1.6

    right = add_textbox(slide, Cm(16.3), Cm(6.0), Cm(11.1), Cm(6.8), fill=WHITE,
                        radius_shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, text="")
    right.text_frame.clear()
    p = right.text_frame.paragraphs[0]
    rr = p.add_run()
    rr.text = "技术栈概览"
    set_east_asian_font(rr)
    rr.font.bold = True
    rr.font.size = Pt(16)
    rr.font.color.rgb = GREEN_DARK
    p.space_after = Pt(10)
    for line in [
        "开发语言：Java 11",
        "界面：XML + Activity 导航",
        "存储：SQLite + DAO 模式",
        "数据准备：JSON 导入 + Gson",
        "中文处理：HanLP、pinyin4j",
        "语音交互：iFlyTek 相关评测/识别模块",
    ]:
        p2 = right.text_frame.add_paragraph()
        p2.text = line
        p2.bullet = True
        p2.space_after = Pt(4)
        for run in p2.runs:
            set_east_asian_font(run)
            run.font.size = Pt(13)
            run.font.color.rgb = TEXT
    add_footer(slide, 5)
    return slide


def make_module_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_slide_title(slide, "5. 核心功能设计", "三类学习任务共同组成 CatLingo 的游戏化训练体系")

    add_card(slide, Cm(1.0), Cm(3.2), Cm(8.6), Cm(6.2), "Character Matching 字符匹配", [
        "根据拼音提示，从 4 个选项中选出正确中文词语。",
        "按词性生成干扰项，避免错误选项过于简单。",
        "目标：强化拼音到汉字/词语的映射识别能力。",
    ], fill=WHITE, title_fill=GREEN_DARK)
    add_card(slide, Cm(10.0), Cm(3.2), Cm(8.6), Cm(6.2), "Pronunciation Quiz 发音测验", [
        "展示目标词与拼音，长按麦克风进行录音评测。",
        "返回总分、流利度、完整度、声调等反馈，再映射为游戏分数。",
        "目标：把发音练习变成可重复的互动任务。",
    ], fill=WHITE, title_fill=GREEN)
    add_card(slide, Cm(19.0), Cm(3.2), Cm(8.6), Cm(6.2), "Word Puzzle 词语谜题", [
        "短句采用排序重组，长句采用挖空填词两种模式。",
        "依赖 HanLP 分词结果与 sentence_words 数据支持题目生成。",
        "目标：训练句序理解与语义组合能力。",
    ], fill=WHITE, title_fill=RGBColor(111, 148, 84))

    bottom = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.0), Cm(10.0), Cm(26.6), Cm(2.1))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = GREEN_LIGHT
    bottom.line.fill.background()
    add_textbox(slide, Cm(1.35), Cm(10.25), Cm(25.8), Cm(1.1),
                "支撑模块：用户登录与注册、难度选择、成绩记录、成就展示、个人资料与学习进度统计。", 16, True, GREEN_DARK)
    fit_image(slide, MASCOT, Cm(23.8), Cm(9.4), width=Cm(2.5), height=Cm(2.5))
    add_footer(slide, 6)
    return slide


def make_database_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_slide_title(slide, "6. 数据库与数据加载机制", "本地持久化保证了无服务器条件下的学习记录与题库管理")

    left = add_textbox(slide, Cm(1.0), Cm(3.0), Cm(8.7), Cm(8.8), fill=WHITE,
                       radius_shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, text="")
    left.text_frame.clear()
    p = left.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "数据库设计要点"
    set_east_asian_font(r)
    r.font.bold = True
    r.font.size = Pt(16)
    r.font.color.rgb = GREEN_DARK
    p.space_after = Pt(8)
    for line in [
        "用户相关：users、achievements、user_achievements",
        "成绩相关：game_scores、game_question_details",
        "资源相关：sentences、sentence_words",
        "题库相关：character_matching、pronunciation_quiz、word_puzzle",
        "采用 DAO + DaoFactory，避免在 Activity 中散落 SQL 代码",
    ]:
        p2 = left.text_frame.add_paragraph()
        p2.text = line
        p2.bullet = True
        p2.space_after = Pt(4)
        for run in p2.runs:
            set_east_asian_font(run)
            run.font.size = Pt(13)
            run.font.color.rgb = TEXT

    pipe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.0), Cm(9.3), Cm(8.7), Cm(2.5))
    pipe.fill.solid()
    pipe.fill.fore_color.rgb = MINT
    pipe.line.fill.background()
    add_textbox(slide, Cm(1.3), Cm(9.55), Cm(8.1), Cm(0.55), "数据管线", 15, True, GREEN_DARK)
    add_textbox(slide, Cm(1.3), Cm(10.2), Cm(8.1), Cm(1.2),
                "JSON 题库变更检测\n→ DataManager 重载\n→ SeedDataLoader 分词/位置更新\n→ SQLite 持久化供游戏调用",
                12, False, TEXT)

    fit_image(slide, ER_DIAGRAM, Cm(10.2), Cm(3.0), width=Cm(17.3), height=Cm(8.8))
    add_footer(slide, 7)
    return slide


def make_results_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_slide_title(slide, "7. 关键实现亮点与阶段性成果", "当前原型已经形成从登录到游戏完成的基本闭环")

    add_card(slide, Cm(1.0), Cm(3.0), Cm(12.7), Cm(8.3), "已完成内容", [
        "登录注册、密码更新、个人资料页面与游戏选择页。",
        "三类小游戏主流程均已有实现分支，支持难度选择与 10 题训练。",
        "成绩记录、题目明细、成就解锁和学习统计已接入本地数据库。",
        "支持 JSON 到 SQLite 的批量导入，便于后续扩充题库。",
    ], fill=WHITE, title_fill=GREEN_DARK)

    add_card(slide, Cm(14.1), Cm(3.0), Cm(6.3), Cm(8.3), "实现亮点", [
        "用单一 GameActivity 统一调度三类游戏模式。",
        "Word Puzzle 结合 HanLP 分词与回退机制提高稳健性。",
        "Pronunciation Quiz 将语音评测结果映射到可读反馈与游戏得分。",
    ], fill=WHITE, title_fill=GREEN)

    status = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(20.8), Cm(3.0), Cm(6.7), Cm(8.3))
    status.fill.solid()
    status.fill.fore_color.rgb = WHITE
    status.line.color.rgb = SAND
    add_textbox(slide, Cm(21.1), Cm(3.25), Cm(6.0), Cm(0.7), "当前状态评估", 16, True, GREEN_DARK)
    items = [
        ("原型完整性", "较高", GREEN),
        ("结构可维护性", "较高", GREEN),
        ("题库扩展性", "较高", GREEN),
        ("语音模块稳定性", "待加强", GOLD),
        ("正式测试与打磨", "待完成", RED),
    ]
    y = 4.25
    for label, value, color in items:
        add_textbox(slide, Cm(21.1), Cm(y), Cm(3.6), Cm(0.6), label, 12, False, TEXT)
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(24.8), Cm(y - 0.04), Cm(1.8), Cm(0.55))
        pill.fill.solid()
        pill.fill.fore_color.rgb = color
        pill.line.fill.background()
        add_textbox(slide, Cm(24.8), Cm(y - 0.02), Cm(1.8), Cm(0.5), value, 10, True, WHITE, PP_ALIGN.CENTER)
        y += 1.0
    add_footer(slide, 8)
    return slide


def make_future_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_slide_title(slide, "8. 局限性、风险与未来工作", "从原型到最终毕业设计仍需补齐的关键部分")

    add_card(slide, Cm(1.0), Cm(3.0), Cm(8.5), Cm(8.2), "当前局限", [
        "语音评测链路仍存在实现叙述与实际模块并行演进的问题，需要统一说明。",
        "系统测试、兼容性测试与性能评估尚未完全展开。",
        "界面细节与用户体验仍有打磨空间，尚未达到最终发布级质量。",
    ], fill=WHITE, title_fill=GREEN_DARK)
    add_card(slide, Cm(10.0), Cm(3.0), Cm(8.5), Cm(8.2), "主要风险", [
        "第三方语音评测可能受口音、环境噪声和设备差异影响。",
        "本地存储包含账户和成绩信息，需要继续优化数据安全处理。",
        "题库规模扩大后，数据更新和题目质量控制会变得更重要。",
    ], fill=WHITE, title_fill=GREEN)
    add_card(slide, Cm(19.0), Cm(3.0), Cm(8.6), Cm(8.2), "未来工作", [
        "统一并稳定发音测评子系统，完善错误处理与重试流程。",
        "补充功能测试、数据库测试、性能测试和用户体验评价。",
        "继续优化 UI 视觉风格，扩展题库、难度策略与学习分析能力。",
    ], fill=WHITE, title_fill=RGBColor(111, 148, 84))
    add_footer(slide, 9)
    return slide


def make_summary_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = CREAM

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    band.fill.solid()
    band.fill.fore_color.rgb = GREEN_DARK
    band.fill.transparency = 0.03
    band.line.fill.background()

    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(2.0), Cm(1.8), Cm(24.0), Cm(11.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.fill.background()

    fit_image(slide, MASCOT, Cm(3.0), Cm(3.1), width=Cm(5.0), height=Cm(5.0))
    add_textbox(slide, Cm(9.0), Cm(3.0), Cm(14.5), Cm(1.4), "总结", 28, True, GREEN_DARK)
    add_textbox(slide, Cm(9.0), Cm(4.4), Cm(14.8), Cm(4.5),
                "本项目围绕“中文学习 + 游戏化交互”构建了一个可运行的 Android 原型。\n\n它已经验证了从用户管理、题库加载到三类小游戏与进度追踪的整体技术路线，并为后续测试完善、语音模块稳定化与最终论文定稿提供了可靠基础。",
                16, False, TEXT)
    thanks = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(9.0), Cm(9.8), Cm(8.0), Cm(1.25))
    thanks.fill.solid()
    thanks.fill.fore_color.rgb = GOLD
    thanks.line.fill.background()
    add_textbox(slide, Cm(9.2), Cm(10.0), Cm(7.6), Cm(0.7), "谢谢老师，敬请指正", 18, True, GREEN_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, Cm(9.0), Cm(11.5), Cm(10.0), Cm(0.8), "Q&A", 20, True, GREEN, PP_ALIGN.LEFT)
    return slide


prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)

make_title_slide()
make_background_slide()
make_objective_slide()
make_requirements_slide()
make_architecture_slide()
make_module_slide()
make_database_slide()
make_results_slide()
make_future_slide()
make_summary_slide()

prs.save(str(OUTPUT))
print(f"Generated: {OUTPUT}")
