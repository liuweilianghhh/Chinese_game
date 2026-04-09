import shutil
from pathlib import Path

from pptx import Presentation


def clear_slides(prs: Presentation) -> None:
    """
    Remove all existing slides from the presentation.

    python-pptx does not expose a public API for this, so we manipulate
    the underlying slide id list and relationships.
    """
    slide_id_list = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slide_ids = list(slide_id_list)
    for sld_id in slide_ids:
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        slide_id_list.remove(sld_id)


def add_section_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    """
    Add a single slide with the standard 'Title and Content' layout.
    """
    # Layout 1 is usually 'Title and Content' in most templates.
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    title_placeholder = slide.shapes.title
    body = None
    # Try to find a body/content placeholder in a robust way
    for ph in slide.placeholders:
        # 1 is usually BODY, but some templates use different indices
        if getattr(ph, "is_placeholder", False) and ph != title_placeholder:
            body = ph
            break

    if title_placeholder is not None:
        title_placeholder.text = title

    if body is not None and hasattr(body, "text_frame") and bullets:
        tf = body.text_frame
        tf.clear()
        # First bullet
        tf.text = bullets[0]
        # Remaining bullets
        for line in bullets[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0


def main() -> None:
    ppt_path = Path(r"d:\学校文件\大四上\Liu_Weiliang_2022213200_MidTerm_PPT.pptx")
    if not ppt_path.exists():
        raise SystemExit(f"PPT file not found: {ppt_path}")

    # Backup original file
    backup_path = ppt_path.with_name(ppt_path.stem + "_backup" + ppt_path.suffix)
    if not backup_path.exists():
        shutil.copy2(ppt_path, backup_path)

    prs = Presentation(ppt_path)

    # Clear all existing slides and rebuild according to the required structure.
    clear_slides(prs)

    # 1. Purpose of the project (problem definition and objectives)
    add_section_slide(
        prs,
        title="Project Purpose / 项目目的",
        bullets=[
            "Problem definition: 缺乏面向初中级学习者的、有趣且系统的中文词汇学习工具",
            "Objectives: 通过小游戏提升词汇量、发音准确性、字词识别与语法应用能力",
            "Target platform: Android 移动端应用",
        ],
    )

    # 2. Background
    add_section_slide(
        prs,
        title="Background / 项目背景",
        bullets=[
            "Existing apps: Duolingo, Baicizhan 等在外语学习中广泛应用",
            "Gap: 中文学习类应用在游戏化、发音评测、学习路径设计方面仍不完善",
            "Motivation: 为中文学习者提供更有趣、可持续的学习体验",
        ],
    )

    # 3. Finished work to date
    add_section_slide(
        prs,
        title="Finished Work to Date / 已完成工作",
        bullets=[
            "Complete UI framework: 登录/注册、游戏选择、个人主页与成就页面",
            "Database & question bank: SQLite 多表结构 + JSON 题库 + HanLP 分词",
            "Two core games: Character Matching & Pronunciation Quiz (集成讯飞 ISE)",
            "Basic reward system: 成就与积分记录、基础统计功能",
        ],
    )

    # 4. Problems and solutions
    add_section_slide(
        prs,
        title="Problems & Solutions / 问题与解决方案",
        bullets=[
            "Platform change: Unity → Native Android; Solution: 采用 Java + XML 原生开发",
            "Integration complexity: 讯飞 ISE WebSocket 与音频流处理; Solution: 封装 IseEvaluator 模块",
            "NLP on mobile: HanLP 大小与性能; Solution: 使用 portable 版 + 预处理写入数据库",
            "Scoring fairness: 非母语者口音差异; Solution: 使用区间映射 + 多维反馈(准确度/流畅度/完整度)",
        ],
    )

    # 5. Next steps
    add_section_slide(
        prs,
        title="Next Steps / 下一步计划",
        bullets=[
            "Implement Word Puzzle game based on existing sentence & segmentation data",
            "Refine achievement system and UI/UX, including result summary screens",
            "Conduct systematic testing (unit, integration, UX) and performance tuning",
            "Prepare final report and user evaluation",
        ],
    )

    prs.save(ppt_path)


if __name__ == "__main__":
    main()

